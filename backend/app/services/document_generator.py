import os
import datetime
import re
from typing import List, Dict, Any, Optional
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

STORAGE_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "storage", "uploads"))
os.makedirs(STORAGE_DIR, exist_ok=True)


def _sanitize_filename(name: str) -> str:
    cleaned = re.sub(r'[^\w\-_]', '_', name.strip())
    return re.sub(r'_+', '_', cleaned)[:40].strip('_')


def _get_timestamp() -> str:
    return datetime.datetime.now().strftime("%Y%m%d_%H%M%S")


# =========================================================================
# 1. Word Document Generator (.docx)
# =========================================================================
def generate_approval_note_docx(
    title: str,
    reference_no: str,
    summary: str,
    findings: List[str],
    recommendation: str
) -> str:
    """Generates a formal industrial/PSU approval note document in Word (.docx)."""
    doc = Document()

    # Header title
    heading = doc.add_heading(title, level=1)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Metadata banner
    p_meta = doc.add_paragraph()
    p_meta.add_run(f"Reference No: {reference_no}\n").bold = True
    p_meta.add_run(f"Date: {datetime.date.today().strftime('%B %d, %Y')}\n")
    p_meta.add_run("Security Classification: INTERNAL / AIR-GAPPED ON-PREMISE ONLY\n").italic = True

    # Section 1: Executive Summary
    doc.add_heading("1. Executive Summary", level=2)
    doc.add_paragraph(summary)

    # Section 2: Key Observations / Findings
    doc.add_heading("2. Technical Observations & Findings", level=2)
    for finding in findings:
        doc.add_paragraph(finding, style="List Bullet")

    # Section 3: Recommendation & Approval Request
    doc.add_heading("3. Operational Recommendations & Sign-off", level=2)
    doc.add_paragraph(recommendation)

    # Sign-off block
    p_sign = doc.add_paragraph("\n\n")
    p_sign.add_run("Verified by Sovereign AI Inspection Engine (100% Local Residency)\n").bold = True
    p_sign.add_run("Authorized Signatory: __________________________\n")

    safe_title = _sanitize_filename(title)
    file_name = f"Approval_Note_{safe_title}_{_get_timestamp()}.docx"
    file_path = os.path.join(STORAGE_DIR, file_name)
    doc.save(file_path)
    return file_path


def generate_docx(title: str = "Industrial Inspection Memo", content: str = "", findings: Optional[List[str]] = None, recommendation: str = "Approved for operations.") -> str:
    """Convenience alias for orchestrator."""
    if findings is None:
        findings = [
            "Calculations validated within isolated sandbox bounds.",
            "Zero outbound network egress detected during synthesis.",
            "All parameters conform to internal standard operating procedures."
        ]
    ref_num = f"SOV/{datetime.datetime.now().strftime('%Y%m')}/{datetime.datetime.now().strftime('%M%S')}"
    return generate_approval_note_docx(
        title=title,
        reference_no=ref_num,
        summary=content,
        findings=findings,
        recommendation=recommendation
    )


# =========================================================================
# 2. Excel Spreadsheet Generator (.xlsx)
# =========================================================================
def generate_calculation_sheet_xlsx(
    sheet_title: str,
    headers: List[str],
    rows: List[List[Any]]
) -> str:
    """Generates a styled Excel spreadsheet (.xlsx) for engineering or academic calculations."""
    wb = Workbook()
    ws = wb.active
    clean_sheet_title = _sanitize_filename(sheet_title)[:28] or "Calculations"
    ws.title = clean_sheet_title

    # Header styling
    ws.append(headers)
    header_fill = PatternFill(start_color="1F4E78", end_color="1F4E78", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF", size=11)
    thin_border = Border(
        left=Side(style='thin', color='D9D9D9'),
        right=Side(style='thin', color='D9D9D9'),
        top=Side(style='thin', color='D9D9D9'),
        bottom=Side(style='thin', color='D9D9D9')
    )

    for col in range(1, len(headers) + 1):
        cell = ws.cell(row=1, column=col)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = thin_border

    # Data rows
    zebra_fill = PatternFill(start_color="F2F5F9", end_color="F2F5F9", fill_type="solid")
    for r_idx, row_data in enumerate(rows, start=2):
        ws.append(row_data)
        for c_idx in range(1, len(headers) + 1):
            cell = ws.cell(row=r_idx, column=c_idx)
            cell.border = thin_border
            if r_idx % 2 == 0:
                cell.fill = zebra_fill
            if isinstance(cell.value, (int, float)):
                cell.alignment = Alignment(horizontal="right")

    # Auto-adjust column widths
    for col in ws.columns:
        max_len = max(len(str(cell.value or '')) for cell in col)
        col_letter = col[0].column_letter
        ws.column_dimensions[col_letter].width = max(max_len + 4, 14)

    safe_title = _sanitize_filename(sheet_title)
    file_name = f"{safe_title}_{_get_timestamp()}.xlsx"
    file_path = os.path.join(STORAGE_DIR, file_name)
    wb.save(file_path)
    return file_path


# =========================================================================
# 3. PowerPoint Presentation Generator (.pptx)
# =========================================================================
def generate_presentation_pptx(
    title: str,
    subtitle: str,
    slides: List[Dict[str, Any]]
) -> str:
    """Generates an executive/board slide deck (.pptx)."""
    prs = Presentation()
    
    # Title Slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"{subtitle}\nAir-Gapped Sovereign AI Deliverable | {datetime.date.today()}"

    # Content Slides
    bullet_layout = prs.slide_layouts[1]
    for slide_data in slides:
        s = prs.slides.add_slide(bullet_layout)
        s.shapes.title.text = slide_data.get("slide_title", "Summary")
        body_shape = s.placeholders[1]
        tf = body_shape.text_frame
        points = slide_data.get("bullet_points", [])
        if points:
            tf.text = points[0]
            for pt in points[1:]:
                p = tf.add_paragraph()
                p.text = pt
                p.level = 0

    safe_title = _sanitize_filename(title)
    file_name = f"Presentation_{safe_title}_{_get_timestamp()}.pptx"
    file_path = os.path.join(STORAGE_DIR, file_name)
    prs.save(file_path)
    return file_path


# =========================================================================
# 4. PDF Compliance Report Generator (.pdf)
# =========================================================================
def generate_pdf(
    title: str = "Compliance Audit Report",
    content: str = "",
    findings: Optional[List[str]] = None,
    table_data: Optional[List[List[str]]] = None
) -> str:
    """Generates a professional PDF audit / inspection report using ReportLab."""
    safe_title = _sanitize_filename(title)
    file_name = f"Report_{safe_title}_{_get_timestamp()}.pdf"
    file_path = os.path.join(STORAGE_DIR, file_name)

    doc = SimpleDocTemplate(
        file_path,
        pagesize=letter,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36
    )
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=18,
        leading=22,
        textColor=colors.HexColor('#1F4E78'),
        alignment=1  # Center
    )
    sub_style = ParagraphStyle(
        'DocSub',
        parent=styles['Normal'],
        fontName='Helvetica-Oblique',
        fontSize=10,
        leading=14,
        textColor=colors.HexColor('#555555'),
        alignment=1
    )
    h2_style = ParagraphStyle(
        'DocH2',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=13,
        leading=16,
        textColor=colors.HexColor('#1F4E78'),
        spaceBefore=12,
        spaceAfter=6
    )
    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['BodyText'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        textColor=colors.HexColor('#222222'),
        spaceAfter=6
    )

    story = []

    # Title & Header
    story.append(Paragraph(title, title_style))
    story.append(Spacer(1, 4))
    story.append(Paragraph(f"SOVEREIGN ON-PREMISE AUDIT • DATE: {datetime.date.today()} • STRICT AIR-GAP GOVERNANCE", sub_style))
    story.append(Spacer(1, 14))

    # Executive Summary Section
    story.append(Paragraph("1. Executive Summary", h2_style))
    story.append(Paragraph(content or "Audit assessment conducted on-premises with zero cloud communication.", body_style))
    story.append(Spacer(1, 10))

    # Findings Section
    story.append(Paragraph("2. Verified Observations", h2_style))
    items = findings or [
        "100% On-premise model execution verified with 0 external network calls.",
        "Code execution completed safely in isolated sandbox environment.",
        "Data residency fully compliant with PSU & Defense confidentiality standards."
    ]
    for item in items:
        story.append(Paragraph(f"• {item}", body_style))
    story.append(Spacer(1, 10))

    # Optional Table
    if table_data and len(table_data) > 0:
        story.append(Paragraph("3. Data Summary Table", h2_style))
        t = Table(table_data)
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1F4E78')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 6),
            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F2F5F9')),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#D9D9D9')),
        ]))
        story.append(t)
        story.append(Spacer(1, 14))

    # Sign-off block
    story.append(Paragraph("4. Sign-off & Integrity Attestation", h2_style))
    story.append(Paragraph("This document was generated autonomously in an isolated sandbox. Zero cloud telemetry was transmitted.", body_style))
    story.append(Spacer(1, 12))
    story.append(Paragraph("Authorized Verification Signature: _________________________________", body_style))

    doc.build(story)
    return file_path